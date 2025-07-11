#!/usr/bin/env python3
"""
url_checker.py

Interactive source selection (single file vs. folder, with optional sub-folder recursion and up to 3 retries)
Per-file URL extraction (PDF, Word, Excel, PowerPoint)
Output CSV named per the spec (_<YYMMDD> suffix), with the proper headers
HTTP checks using a disguised, randomized User-Agent and cloudscraper to handle Cloudflare/CAPTCHA
Human-like behavior via random pauses
Dependency checks that prompt before installing any missing package

Adds support for .txt, .csv, and OpenDocument formats (.odt, .ods, .odp)
to the existing PDF, Word, PPTX, and XLSX extractors.

"""

import os
import sys
import re
import csv
import time
import random
import zipfile
from datetime import datetime
from pathlib import Path

# 1) DEPENDENCY MANAGEMENT
REQUIRED_PACKAGES = {
    'cloudscraper': 'cloudscraper',
    'PyPDF2': 'PyPDF2',
    'docx': 'python-docx',
    'pptx': 'python-pptx',
    'openpyxl': 'openpyxl',
}

def check_and_install(pkg_name, install_name=None):
    install_name = install_name or pkg_name
    try:
        __import__(pkg_name)
    except ImportError:
        resp = input(f"Package '{install_name}' is required. Install now? "
                     "(may modify your environment) [y/N]: ").strip().lower()
        if resp == 'y':
            import subprocess
            subprocess.check_call([sys.executable, '-m', 'pip', 'install', install_name])
        else:
            print(f"Cannot continue without '{install_name}'. Exiting.")
            sys.exit(1)

for module, pkg in REQUIRED_PACKAGES.items():
    check_and_install(module, pkg)

import cloudscraper
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl

# 2) GLOBALS
DATE_SUFFIX = datetime.now().strftime("%y%m%d")
RESPONSE_CODES = {
    100: "Continue", 101: "Switching Protocols",
    200: "OK", 201: "Created", 202: "Accepted", 204: "No Content",
    301: "Moved Permanently", 302: "Found", 304: "Not Modified",
    400: "Bad Request", 401: "Unauthorized", 403: "Forbidden", 404: "Not Found",
    500: "Internal Server Error", 502: "Bad Gateway",
    503: "Service Unavailable", 504: "Gateway Timeout"
}

USER_AGENTS = [
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
      "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/114.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
      "AppleWebKit/605.1.15 (KHTML, like Gecko) Version/15.1 Safari/605.1.15",
    "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 "
      "(KHTML, like Gecko) Ubuntu Chromium/113.0.0.0 Chrome/113.0.0.0 Safari/537.36",
]

URL_REGEX = re.compile(r'https?://[^\s\)\"\'>]+')

def get_session():
    scraper = cloudscraper.create_scraper()
    ua = random.choice(USER_AGENTS)
    scraper.headers.update({
        'User-Agent': ua,
        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
        'Accept-Language': 'en-US,en;q=0.5',
        'Connection': 'keep-alive',
    })
    return scraper

def human_sleep(a=1.0, b=3.0):
    time.sleep(random.uniform(a, b))

# 3) VALIDATION LOGIC (up to 3 tries)
def validate_source(check_all_files: bool):
    tries, max_tries = 0, 3
    prompt = ("Path to folder" if check_all_files else "Path to file") + ": "
    while tries < max_tries:
        src = input(prompt).strip()
        p = Path(src)
        if check_all_files and p.is_dir():
            return p
        if not check_all_files and p.is_file():
            return p
        tries += 1
        print(f"Invalid path. {max_tries - tries} tries remaining.")
        if tries < max_tries:
            again = input("Try again? [Y/n]: ").strip().lower()
            if again in ('n', 'no'):
                print("Exiting.")
                sys.exit(1)
    print(f"Reached {max_tries} invalid attempts. Exiting.")
    sys.exit(1)

# 4) URL EXTRACTION PER FILETYPE
def extract_urls_from_pdf(path: Path):
    text = ""
    for page in PdfReader(str(path)).pages:
        text += page.extract_text() or ""
    return set(URL_REGEX.findall(text))

def extract_urls_from_docx(path: Path):
    doc = Document(str(path))
    full = "\n".join(p.text for p in doc.paragraphs)
    return set(URL_REGEX.findall(full))

def extract_urls_from_pptx(path: Path):
    prs = Presentation(str(path))
    found = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                found |= set(URL_REGEX.findall(shape.text))
            if shape.element.hyperlink:
                found.add(shape.element.hyperlink.target)
    return found

def extract_urls_from_xlsx(path: Path):
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    urls = set()
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if isinstance(cell, str):
                    urls |= set(URL_REGEX.findall(cell))
    return urls

def extract_urls_from_txt(path: Path):
    try:
        text = path.read_text(encoding='utf-8', errors='ignore')
        return set(URL_REGEX.findall(text))
    except Exception as e:
        print(f"Failed to read TXT {path}: {e}")
        return set()

def extract_urls_from_csv(path: Path):
    urls = set()
    try:
        with open(path, newline='', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for row in reader:
                for cell in row:
                    if isinstance(cell, str):
                        urls |= set(URL_REGEX.findall(cell))
    except Exception as e:
        print(f"Failed to read CSV {path}: {e}")
    return urls

def extract_urls_from_odf(path: Path):
    urls = set()
    try:
        with zipfile.ZipFile(str(path)) as z:
            if 'content.xml' in z.namelist():
                xml = z.read('content.xml').decode('utf-8', errors='ignore')
                urls |= set(URL_REGEX.findall(xml))
    except Exception as e:
        print(f"Failed to extract ODF {path}: {e}")
    return urls

# Map extensions to extractor functions
EXTRACTORS = {
    '.pdf': extract_urls_from_pdf,
    '.docx': extract_urls_from_docx,
    '.pptx': extract_urls_from_pptx,
    '.xlsx': extract_urls_from_xlsx,
    '.txt': extract_urls_from_txt,
    '.csv': extract_urls_from_csv,
    '.odt': extract_urls_from_odf,
    '.ods': extract_urls_from_odf,
    '.odp': extract_urls_from_odf,
}

# 5) HTTP CHECKING
def check_url(session, url, max_retries=3):
    for attempt in range(1, max_retries+1):
        try:
            r = session.get(url, timeout=10)
            code = r.status_code
            if code in (429, 503) or 'captcha' in (r.text or '').lower():
                wait = random.uniform(5, 10)
                print(f"Encountered {code} on {url}, waiting {wait:.1f}s before retry...")
                time.sleep(wait)
                continue
            return code, RESPONSE_CODES.get(code, "Unknown")
        except Exception as e:
            print(f"Error checking {url}: {e}")
            time.sleep(random.uniform(2,4))
    return None, "Failed after retries"

# 6) PROCESSING
def process_file(path: Path, writer, session):
    extractor = EXTRACTORS.get(path.suffix.lower())
    if not extractor:
        return
    urls = extractor(path)
    for url in urls:
        human_sleep()
        code, comment = check_url(session, url)
        writer.writerow({
            'URL': url,
            'response_code': code or '',
            'comment': comment or ''
        })

def process_single(path: Path):
    out_name = f"{path.stem}_url_check_results_{DATE_SUFFIX}.csv"
    out_path = path.parent / out_name
    with open(out_path, 'w', newline='', encoding='utf-8') as csvfile:
        fieldnames = ['URL','response_code','comment']
        writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
        writer.writeheader()
        session = get_session()
        process_file(path, writer, session)
    print(f"Results written to {out_path}")

def process_folder(folder: Path, include_subs: bool):
    out_name = f"{folder.name}_url_check_results_{DATE_SUFFIX}.csv"
    out_path = folder / out_name
    with open(out_path, 'w', newline='', encoding='utf-8') as csvfile:
        fieldnames = ['folder_path','file_name','URL','response_code','comment']
        writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
        writer.writeheader()

        all_files = (
            list(folder.rglob('*')) if include_subs else list(folder.iterdir())
        )
        files = [p for p in all_files if p.suffix.lower() in EXTRACTORS]
        total = len(files)
        session = get_session()
        for idx, f in enumerate(files, 1):
            extractor = EXTRACTORS[f.suffix.lower()]
            urls = extractor(f)
            for url in urls:
                human_sleep()
                code, comment = check_url(session, url)
                writer.writerow({
                    'folder_path': str(folder),
                    'file_name': f.name,
                    'URL': url,
                    'response_code': code or '',
                    'comment': comment or ''
                })
            if idx == 1 or (idx/total)*100 % 5 < (100/total):
                print(f"Processed {idx}/{total} files ({(idx/total)*100:.0f}%)")

    print(f"Results written to {out_path}")

# 7) MAIN
def main():
    print("=== URL Checker ===")
    # NEW: inform the user of all supported file extensions
    supported = sorted(EXTRACTORS.keys())
    print("Supported file formats:", ", ".join(supported))

    all_files = input("Check ALL files in a folder? [y/N]: ").strip().lower() in ('y','yes')
    sub_folders = False
    if all_files:
        sub_folders = input("Include ALL sub-folders? [y/N]: ").strip().lower() in ('y','yes')
    source = validate_source(all_files)
    if all_files:
        process_folder(source, sub_folders)
    else:
        process_single(source)

if __name__ == "__main__":
    main()
